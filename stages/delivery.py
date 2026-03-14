"""
Этап 3: Delivery — генерирует PPTX.
Consultant-grade дизайн: чистая геометрия, action titles, data-forward.
Все цвета, шрифты, размеры берутся из session["brand"].
"""
from __future__ import annotations
import io, tempfile, logging
import matplotlib; matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from stages.layout_registry import truncate_content

log = logging.getLogger(__name__)

# ── Canvas & Grid ────────────────────────────────────────────────────────────
_W = Inches(13.33)
_H = Inches(7.5)

MARGIN     = 1.0
CONTENT_W  = 11.33
HEADER_Y   = 0.4
BODY_TOP   = 1.8
BODY_BOTTOM = 6.8
FOOTER_Y   = 7.0


# ── Logo Cache ───────────────────────────────────────────────────────────────
class LogoCache:
    """Downloads logo once, caches bytes + real aspect ratio."""

    def __init__(self):
        self._data: bytes | None = None
        self._aspect: float = 0.4  # fallback w/h ratio
        self._loaded = False

    def get(self, brand) -> tuple[bytes | None, float]:
        if self._loaded:
            return self._data, self._aspect
        self._loaded = True
        if not brand.logo.url:
            return None, self._aspect
        try:
            import urllib.request
            self._data = urllib.request.urlopen(brand.logo.url, timeout=5).read()
            from PIL import Image as PILImage
            img = PILImage.open(io.BytesIO(self._data))
            pw, ph = img.size
            self._aspect = ph / pw if pw else 0.4
        except Exception:
            pass
        return self._data, self._aspect


# ── Overflow Guard ───────────────────────────────────────────────────────────
class OverflowGuard:
    """Tracks current y position and checks if an element fits."""

    def __init__(self, y_start: float, y_limit: float):
        self.y = y_start
        self.y_limit = y_limit

    def can_fit(self, h: float) -> bool:
        return self.y + h <= self.y_limit

    def advance(self, h: float):
        self.y += h


# ── Utilities ────────────────────────────────────────────────────────────────
def _rgb_to_hex(color: RGBColor) -> str:
    return f"#{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def _rect(sl, x, y, w, h, color: RGBColor):
    sh = sl.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()


def _rounded_rect(sl, x, y, w, h, fill_color: RGBColor):
    sh = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    sh.line.fill.background()
    return sh


def _accent_rule(sl, x, y, w, color: RGBColor):
    _rect(sl, x, y, w, 0.03, color)


_FALLBACK_FONTS = ["Calibri", "Arial"]


def _text(sl, text, x, y, w, h, size=18, bold=False, color: RGBColor = None,
          align=PP_ALIGN.LEFT, font=None, brand=None,
          line_spacing=1.15, space_before=0, space_after=0):
    font_name = font or (brand.typography.font_body if brand else _FALLBACK_FONTS[0])
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or (brand.colors.text_dark if brand else RGBColor(0, 0, 0))
    run.font.name = font_name
    return tb


# ── Slide components ─────────────────────────────────────────────────────────
def _slide_background(sl, color: RGBColor):
    """Full-slide rect as background (compatible with Keynote)."""
    _rect(sl, 0, 0, 13.33, 7.5, color)


def _modern_header(sl, title, brand):
    """Action title + accent rule (replaces full-width color bar)."""
    _text(sl, title, MARGIN, HEADER_Y, CONTENT_W, 1.2,
          size=brand.typography.size_heading, bold=True,
          color=brand.colors.text_dark,
          font=brand.typography.font_heading, brand=brand,
          line_spacing=1.1)
    _accent_rule(sl, MARGIN, 1.55, 2.0, brand.colors.accent)


def _modern_footer(sl, slide_num, total, brand):
    """Thin accent line + company name left + slide number right."""
    _accent_rule(sl, MARGIN, FOOTER_Y, CONTENT_W, brand.colors.accent)
    if brand.slide_defaults.footer_text:
        _text(sl, brand.slide_defaults.footer_text, MARGIN, FOOTER_Y + 0.08, 6, 0.3,
              size=brand.typography.size_caption,
              color=brand.colors.text_muted, brand=brand)
    if brand.slide_defaults.slide_numbers:
        _text(sl, f"{slide_num}/{total}", 11.5, FOOTER_Y + 0.08, 0.83, 0.3,
              size=brand.typography.size_caption,
              color=brand.colors.text_muted, align=PP_ALIGN.RIGHT, brand=brand)


def _cached_logo(sl, brand, logo_cache: LogoCache):
    data, aspect = logo_cache.get(brand)
    if not data:
        return
    w = brand.logo.width_inches
    h = w * aspect
    pos = brand.logo.position
    x = 13.33 - w - 0.2 if "right" in pos else 0.2
    y = 7.5 - h - 0.05 if "bottom" in pos else 0.05
    sl.shapes.add_picture(
        io.BytesIO(data), Inches(x), Inches(y),
        width=Inches(w), height=Inches(h),
    )


# ── Slide renderers ──────────────────────────────────────────────────────────

def _slide_title(sl, s, brand, logo_cache):
    """Dark full-bleed, LEFT-aligned title slide."""
    _slide_background(sl, brand.colors.bg_dark)

    _text(sl, s.get("title", ""), MARGIN, 2.2, CONTENT_W, 1.6,
          size=brand.typography.size_title, bold=True,
          color=brand.colors.bg_light, align=PP_ALIGN.LEFT,
          font=brand.typography.font_heading, brand=brand,
          line_spacing=1.05)

    _accent_rule(sl, MARGIN, 3.9, 3.0, brand.colors.accent)

    sub = s.get("subtitle", "")
    if sub:
        _text(sl, sub, MARGIN, 4.15, CONTENT_W, 0.8,
              size=brand.typography.size_section,
              color=brand.colors.text_muted, brand=brand)

    # Date
    from datetime import date
    lang = brand.language
    today = date.today()
    if lang == "ru":
        _MONTHS_RU = [
            "", "января", "февраля", "марта", "апреля", "мая", "июня",
            "июля", "августа", "сентября", "октября", "ноября", "декабря",
        ]
        date_str = f"{today.day} {_MONTHS_RU[today.month]} {today.year}"
    else:
        date_str = today.strftime("%B %d, %Y")
    _text(sl, date_str, MARGIN, 6.4, 4, 0.4,
          size=brand.typography.size_caption,
          color=brand.colors.text_muted, brand=brand)

    _cached_logo(sl, brand, logo_cache)


def _slide_content(sl, s, brand, logo_cache, slide_num, total):
    """Content slide with bullets and highlights, overflow-protected."""
    _slide_background(sl, brand.colors.bg_light)
    _modern_header(sl, s.get("title", ""), brand)

    guard = OverflowGuard(BODY_TOP, BODY_BOTTOM)
    items = s.get("content", [])[:6]  # max 6 items

    for item in items:
        t = item.get("type", "bullet")
        txt = item.get("text", "")
        if t == "highlight":
            if not guard.can_fit(0.85):
                break
            _rect(sl, MARGIN, guard.y, 0.06, 0.65, brand.colors.accent)
            _text(sl, txt, 1.3, guard.y, 10.7, 0.7,
                  size=brand.typography.size_body, bold=True,
                  color=brand.colors.accent, brand=brand)
            guard.advance(0.85)
        else:
            if not guard.can_fit(0.75):
                break
            _text(sl, f"\u2022  {txt}", 1.3, guard.y, 10.7, 0.7,
                  size=brand.typography.size_body,
                  color=brand.colors.text_dark, brand=brand)
            guard.advance(0.75)

    if s.get("speaker_notes"):
        sl.notes_slide.notes_text_frame.text = s["speaker_notes"]

    _modern_footer(sl, slide_num, total, brand)
    _cached_logo(sl, brand, logo_cache)


def _slide_chart(sl, s, brand, charts_data, logo_cache, slide_num, total):
    """Chart slide with proportionally scaled chart image."""
    _slide_background(sl, brand.colors.bg_light)
    _modern_header(sl, s.get("title", ""), brand)

    idx = s.get("chart_ref", 0)
    cd = charts_data[idx] if idx < len(charts_data) else None
    if cd:
        img = _render_chart(cd, brand)
        from PIL import Image as PILImage
        pi = PILImage.open(io.BytesIO(img))
        pw, ph = pi.size

        # Max area: (MARGIN, BODY_TOP, CONTENT_W, 4.8)
        max_w, max_h = CONTENT_W, 4.8
        scale = min(max_w / (pw / 192), max_h / (ph / 192))
        iw = (pw / 192) * scale
        ih = (ph / 192) * scale
        ix = MARGIN + (max_w - iw) / 2
        iy = BODY_TOP + (max_h - ih) / 2

        sl.shapes.add_picture(
            io.BytesIO(img), Inches(ix), Inches(iy),
            width=Inches(iw), height=Inches(ih),
        )

    _modern_footer(sl, slide_num, total, brand)
    _cached_logo(sl, brand, logo_cache)


def _slide_two_col(sl, s, brand, logo_cache, slide_num, total):
    """Two-column slide with vertical accent divider."""
    _slide_background(sl, brand.colors.bg_light)
    _modern_header(sl, s.get("title", ""), brand)

    # Vertical divider
    _rect(sl, 6.5, 2.0, 0.015, 4.5, brand.colors.accent)

    for col, x, w in [
        (s.get("left", {}), MARGIN, 5.0),
        (s.get("right", {}), 7.0, 5.33),
    ]:
        heading = col.get("heading", "")
        if heading:
            _text(sl, heading, x, BODY_TOP, w, 0.5,
                  size=brand.typography.size_section, bold=True,
                  color=brand.colors.text_dark,
                  font=brand.typography.font_heading, brand=brand)
            _accent_rule(sl, x, BODY_TOP + 0.5, min(w, 2.0), brand.colors.accent)

        guard = OverflowGuard(BODY_TOP + 0.7, BODY_BOTTOM)
        for item in col.get("items", [])[:5]:
            if not guard.can_fit(0.65):
                break
            _text(sl, f"\u2022  {item}", x, guard.y, w, 0.6,
                  size=brand.typography.size_body,
                  color=brand.colors.text_dark, brand=brand)
            guard.advance(0.65)

    _modern_footer(sl, slide_num, total, brand)
    _cached_logo(sl, brand, logo_cache)


def _slide_stats(sl, s, brand, logo_cache, slide_num, total):
    """Stats slide with rounded-rect cards and accent top-bar."""
    _slide_background(sl, brand.colors.bg_light)
    _modern_header(sl, s.get("title", ""), brand)

    stats = s.get("stats", [])[:4]
    n = max(len(stats), 1)
    gap = 0.4
    total_gap = gap * (n - 1)
    card_w = (CONTENT_W - total_gap) / n
    card_h = 3.6
    card_y = BODY_TOP + 0.2

    for i, st in enumerate(stats):
        cx = MARGIN + i * (card_w + gap)

        # Card background
        _rounded_rect(sl, cx, card_y, card_w, card_h, RGBColor(0xFF, 0xFF, 0xFF))
        # Accent top-bar
        _rect(sl, cx + 0.1, card_y, card_w - 0.2, 0.06, brand.colors.accent)

        # Hero value — dynamic font size
        val = st.get("value", "")
        if len(val) > 10:
            val_size = 24
        elif len(val) > 6:
            val_size = 32
        else:
            val_size = 40

        _text(sl, val, cx, card_y + 0.5, card_w, 1.2,
              size=val_size, bold=True, color=brand.colors.primary,
              align=PP_ALIGN.CENTER, font=brand.typography.font_heading,
              brand=brand)

        # Trend arrow
        trend = st.get("trend", "")
        if trend:
            if any(c in trend for c in ["+", "\u2191"]):
                tc = brand.colors.success
                arrow = "\u2191 "
            elif any(c in trend for c in ["-", "\u2193"]):
                tc = brand.colors.danger
                arrow = "\u2193 "
            else:
                tc = brand.colors.text_muted
                arrow = "\u2192 "
            # Prepend arrow if not already there
            if trend[0] not in "\u2191\u2193\u2192":
                trend = arrow + trend
            _text(sl, trend, cx, card_y + 1.7, card_w, 0.5,
                  size=18, bold=True, color=tc,
                  align=PP_ALIGN.CENTER, brand=brand)

        # Label
        _text(sl, st.get("label", ""), cx, card_y + 2.4, card_w, 0.8,
              size=brand.typography.size_caption,
              color=brand.colors.text_muted,
              align=PP_ALIGN.CENTER, brand=brand)

    _modern_footer(sl, slide_num, total, brand)
    _cached_logo(sl, brand, logo_cache)


def _slide_closing(sl, s, brand, logo_cache):
    """Dark full-bleed closing, LEFT-aligned."""
    _slide_background(sl, brand.colors.bg_dark)

    _text(sl, s.get("title", "Спасибо!"), MARGIN, 1.5, CONTENT_W, 1.4,
          size=brand.typography.size_title, bold=True,
          color=brand.colors.bg_light, align=PP_ALIGN.LEFT,
          font=brand.typography.font_heading, brand=brand)

    _accent_rule(sl, MARGIN, 3.0, 3.0, brand.colors.accent)

    guard = OverflowGuard(3.3, 6.2)
    items = s.get("content", [])[:4]
    for item in items:
        txt = item.get("text", "")
        if not guard.can_fit(0.7):
            break
        _text(sl, f"\u2014  {txt}", MARGIN, guard.y, CONTENT_W, 0.65,
              size=brand.typography.size_body,
              color=brand.colors.accent, align=PP_ALIGN.LEFT, brand=brand)
        guard.advance(0.7)

    # Footer text (no number)
    if brand.slide_defaults.footer_text:
        _text(sl, brand.slide_defaults.footer_text, MARGIN, FOOTER_Y + 0.08, 6, 0.3,
              size=brand.typography.size_caption,
              color=brand.colors.text_muted, brand=brand)

    _cached_logo(sl, brand, logo_cache)


def _slide_section(sl, s, brand):
    """Section divider — dark full-bleed, number + title."""
    _slide_background(sl, brand.colors.bg_dark)

    num = s.get("section_number", "")
    if num:
        _text(sl, num, MARGIN, 1.5, 4, 1.8,
              size=72, bold=True,
              color=brand.colors.accent,
              font=brand.typography.font_heading, brand=brand)

    _accent_rule(sl, MARGIN, 3.6, 2.0, brand.colors.accent)

    _text(sl, s.get("title", ""), MARGIN, 3.9, CONTENT_W, 1.2,
          size=brand.typography.size_title, bold=True,
          color=brand.colors.bg_light,
          font=brand.typography.font_heading, brand=brand,
          line_spacing=1.05)

    sub = s.get("subtitle", "")
    if sub:
        _text(sl, sub, MARGIN, 5.4, CONTENT_W, 0.8,
              size=brand.typography.size_section,
              color=brand.colors.text_muted, brand=brand)


# ── Chart rendering ──────────────────────────────────────────────────────────
def _render_chart(cd: dict, brand) -> bytes:
    kind = cd.get("kind", "bar")
    labels = cd.get("labels", [])
    series = cd.get("series", [])
    colors = brand.chart_colors(len(series) or 1)

    bg_hex = _rgb_to_hex(brand.colors.bg_light)
    muted_hex = _rgb_to_hex(brand.colors.text_muted)
    grid_color = "#E2E8F0"

    try:
        font_body = brand.typography.font_body
        plt.rcParams["font.family"] = font_body
    except Exception:
        pass

    fig, ax = plt.subplots(figsize=(11, 5.0))
    fig.patch.set_facecolor(bg_hex)
    ax.set_facecolor(bg_hex)
    for sp in ["top", "right"]:
        ax.spines[sp].set_visible(False)
    for sp in ["left", "bottom"]:
        ax.spines[sp].set_color(grid_color)
    ax.tick_params(colors=muted_hex, labelsize=11)
    ax.yaxis.grid(True, color=grid_color, linewidth=0.7)
    ax.set_axisbelow(True)

    if kind == "bar":
        x = np.arange(len(labels))
        bw = 0.7 / max(len(series), 1)
        for i, s in enumerate(series):
            off = (i - (len(series) - 1) / 2) * bw
            bars = ax.bar(
                x + off, s["values"], width=bw * 0.88,
                color=colors[i], label=s.get("name", ""),
            )
            ax.bar_label(bars, fmt="%.0f", fontsize=9, color=muted_hex, padding=3)
        ax.set_xticks(x)
        ax.set_xticklabels(labels, fontsize=11)

    elif kind == "line":
        for i, s in enumerate(series):
            ax.plot(
                labels, s["values"], color=colors[i],
                linewidth=2.5, marker="o", markersize=5,
                label=s.get("name", ""),
            )

    elif kind in ("pie", "doughnut"):
        vals = series[0]["values"] if series else []
        wp = {"width": 0.5} if kind == "doughnut" else {}
        ax.pie(
            vals, labels=labels, colors=colors,
            autopct="%1.0f%%", startangle=90, wedgeprops=wp,
            textprops={"fontsize": 11, "color": _rgb_to_hex(brand.colors.text_dark)},
        )

    if cd.get("title"):
        ax.set_title(
            cd["title"], fontsize=14, fontweight="bold",
            color=brand.colors.primary_hex, pad=12,
        )
    if len(series) > 1:
        ax.legend(fontsize=11, framealpha=0)

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=192, bbox_inches="tight", facecolor=bg_hex)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


# ── PPTX builder ─────────────────────────────────────────────────────────────
def _build_pptx(slides: list, charts_data: list, path: str, brand):
    # Truncate content to fit layout constraints before rendering
    slides = [truncate_content(s) for s in slides]

    prs = Presentation()
    prs.slide_width = _W
    prs.slide_height = _H
    # Metadata for portability (ClaWic rule 6)
    prs.core_properties.title = brand.company_name
    prs.core_properties.subject = "Generated by PresMaker"
    blank = prs.slide_layouts[6]
    logo_cache = LogoCache()
    total = len(slides)

    for idx, s in enumerate(slides):
        sl = prs.slides.add_slide(blank)
        t = s.get("type", "content")
        num = idx + 1

        if t == "title":
            _slide_title(sl, s, brand, logo_cache)
        elif t == "section":
            _slide_section(sl, s, brand)
        elif t == "chart":
            _slide_chart(sl, s, brand, charts_data, logo_cache, num, total)
        elif t == "two_column":
            _slide_two_col(sl, s, brand, logo_cache, num, total)
        elif t == "stats":
            _slide_stats(sl, s, brand, logo_cache, num, total)
        elif t == "closing":
            _slide_closing(sl, s, brand, logo_cache)
        else:
            _slide_content(sl, s, brand, logo_cache, num, total)

    prs.save(path)


# ── Stage class ──────────────────────────────────────────────────────────────
class DeliveryBuildStage:

    async def run(self, session) -> tuple[str, dict]:
        brand = session["brand"]
        slides = session["filled_slides"]
        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        tmp.close()
        _build_pptx(
            slides,
            session["research_data"].get("data_for_charts", []),
            tmp.name,
            brand,
        )
        return tmp.name, {
            "title": session["research_data"].get("topic", "Презентация"),
            "slides": len(slides),
        }
